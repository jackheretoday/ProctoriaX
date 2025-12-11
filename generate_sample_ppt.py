"""
Generate a sample PowerPoint file with properly formatted questions
Run this script to create a test PPTX file that matches the expected format
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_sample_questions_ppt():
    """Create a sample PowerPoint with 5 properly formatted questions"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Sample questions
    questions = [
        {
            'number': 1,
            'text': 'What is Python?',
            'options': {
                'A': 'A snake',
                'B': 'A programming language',
                'C': 'A game',
                'D': 'None of the above'
            },
            'answer': 'B',
            'explanation': 'Python is a high-level programming language used for web development, data science, and automation.'
        },
        {
            'number': 2,
            'text': 'Which of the following is a valid Python data type?',
            'options': {
                'A': 'Integer',
                'B': 'String',
                'C': 'List',
                'D': 'All of the above'
            },
            'answer': 'D',
            'explanation': 'Python supports multiple data types including integers, strings, lists, tuples, and dictionaries.'
        },
        {
            'number': 3,
            'text': 'What does HTML stand for?',
            'options': {
                'A': 'Hyper Text Markup Language',
                'B': 'High Tech Modern Language',
                'C': 'Home Tool Markup Language',
                'D': 'Hyperlinks and Text Markup Language'
            },
            'answer': 'A',
            'explanation': 'HTML stands for Hyper Text Markup Language and is the standard markup language for creating web pages.'
        },
        {
            'number': 4,
            'text': 'What is the result of 2 + 2 in Python?',
            'options': {
                'A': '3',
                'B': '4',
                'C': '5',
                'D': '22'
            },
            'answer': 'B',
            'explanation': 'In Python, the + operator performs addition for numbers, so 2 + 2 equals 4.'
        },
        {
            'number': 5,
            'text': 'Which keyword is used to define a function in Python?',
            'options': {
                'A': 'function',
                'B': 'def',
                'C': 'func',
                'D': 'define'
            },
            'answer': 'B',
            'explanation': 'The "def" keyword is used to define a function in Python.'
        }
    ]
    
    # Create a slide for each question
    for q in questions:
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add text box with question content
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Build question text
        question_text = f"Question {q['number']}: {q['text']}\n"
        question_text += f"A) {q['options']['A']}\n"
        question_text += f"B) {q['options']['B']}\n"
        question_text += f"C) {q['options']['C']}\n"
        question_text += f"D) {q['options']['D']}\n"
        question_text += f"Answer: {q['answer']}\n"
        question_text += f"Explanation: {q['explanation']}"
        
        # Add text
        p = text_frame.paragraphs[0]
        p.text = question_text
        p.font.size = Pt(18)
        p.font.name = 'Arial'
    
    # Save presentation
    output_file = 'sample_questions.pptx'
    prs.save(output_file)
    print(f"✅ Sample PowerPoint created: {output_file}")
    print(f"📊 Total questions: {len(questions)}")
    print(f"\nYou can now upload this file via:")
    print("  Teacher Dashboard → Upload Questions → Select Test → Upload File")
    
    return output_file

if __name__ == '__main__':
    try:
        create_sample_questions_ppt()
    except ImportError as e:
        print("❌ Error: python-pptx library not installed")
        print("Run: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating PowerPoint: {str(e)}")
