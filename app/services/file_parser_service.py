"""
File Parser Service
Parses Word and PowerPoint files to extract questions and terms & conditions
"""
import re
from docx import Document
from pptx import Presentation
from app.utils.exceptions import ValidationError


class FileParserService:
    """Service for parsing uploaded files"""
    
    @staticmethod
    def parse_word_document(file_path):
        """
        Parse Word document (.docx) to extract questions
        
        Args:
            file_path: Path to .docx file
            
        Returns:
            List of question dictionaries
            
        Raises:
            ValidationError: If file format is invalid
        """
        try:
            doc = Document(file_path)
            full_text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
            
            questions = FileParserService._parse_questions_from_text(full_text)
            
            if not questions:
                raise ValidationError("No valid questions found in the document")
            
            return questions
            
        except Exception as e:
            if isinstance(e, ValidationError):
                raise
            raise ValidationError(f"Error parsing Word document: {str(e)}")
    
    @staticmethod
    def parse_powerpoint(file_path):
        """
        Parse PowerPoint (.pptx) to extract questions (one question per slide)
        
        Args:
            file_path: Path to .pptx file
            
        Returns:
            List of question dictionaries
            
        Raises:
            ValidationError: If file format is invalid
        """
        try:
            prs = Presentation(file_path)
            questions = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                # Extract all text from slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if not slide_text:
                    continue
                
                # Join all text from the slide
                full_slide_text = '\n'.join(slide_text)
                
                # Parse the slide as a single question
                parsed_questions = FileParserService._parse_questions_from_text(full_slide_text)
                
                if parsed_questions:
                    # Take only the first question from each slide
                    questions.append(parsed_questions[0])
            
            if not questions:
                raise ValidationError("No valid questions found in the PowerPoint. Ensure each slide has: Question text, Options A-D, and Answer line.")
            
            return questions
            
        except Exception as e:
            if isinstance(e, ValidationError):
                raise
            raise ValidationError(f"Error parsing PowerPoint: {str(e)}")
    
    @staticmethod
    def _parse_questions_from_text(text):
        """
        Parse questions from text using regex patterns
        
        Args:
            text: Full text containing questions
            
        Returns:
            List of question dictionaries
        """
        questions = []
        
        # Try to split by "Question X:" pattern first (for multi-line format)
        question_blocks = re.split(r'(?=Question\s+\d+[:\.\)])', text, flags=re.MULTILINE | re.IGNORECASE)
        
        # If no clear "Question X:" separators, try other patterns
        if len(question_blocks) == 1:
            question_blocks = re.split(r'(?=Q\d+[:\.\)])', text, flags=re.MULTILINE | re.IGNORECASE)
        
        if len(question_blocks) == 1:
            # Try splitting by numbered lines
            question_blocks = re.split(r'(?=^\d+[\.\)])', text, flags=re.MULTILINE | re.IGNORECASE)
        
        for block in question_blocks:
            if not block.strip():
                continue
            
            try:
                question_dict = FileParserService._parse_single_question(block)
                if question_dict:
                    questions.append(question_dict)
            except:
                continue
        
        return questions
    
    @staticmethod
    def _parse_single_question(question_text):
        """
        Parse a single question block into structured data
        
        Expected format:
        Question 1: Who developed Python programming language?
        A) Dennis Ritchie
        B) Bjarne Stroustrup
        C) Guido van Rossum
        D) James Gosling
        Answer: C
        Explanation: Python was developed by Guido van Rossum in 1991.
        
        Also handles single-line format:
        Question 1: Who developed Python programming language? A) Dennis Ritchie B) Bjarne Stroustrup C) Guido van Rossum D) James Gosling Answer: C Explanation: Python was developed by Guido van Rossum in 1991.
        
        Args:
            question_text: Text containing a single question
            
        Returns:
            Question dictionary or None
        """
        try:
            # Extract question number and text - more flexible pattern
            # Handles "Question 1:", "Question 1", "Q1:", "1." etc.
            question_num_match = re.search(r'(?:Question\s+|Q)?(\d+)[\s\:\.]+(.*?)(?=\s+[A-D][\)\.]|$)', 
                                          question_text, re.DOTALL | re.IGNORECASE)
            
            if not question_num_match:
                return None
            
            question_number = int(question_num_match.group(1))
            question_text_content = question_num_match.group(2).strip()
            
            # Remove "Question X:" prefix if it exists in the question text
            question_text_content = re.sub(r'^(?:Question\s+\d+|Q\d+)[\s\:\.]\)?\s*', '', 
                                           question_text_content, flags=re.IGNORECASE).strip()
            
            # Extract options (A, B, C, D) - handle both single-line and multi-line formats
            options = {}
            
            # First try multi-line format (each option on new line)
            multiline_option_pattern = r'^([A-D])[\)\.]\s*(.+?)(?=\n[A-D][\)\.]|\nAnswer|\nExplanation|\Z)'
            multiline_matches = re.finditer(multiline_option_pattern, question_text, re.MULTILINE | re.DOTALL | re.IGNORECASE)
            
            for match in multiline_matches:
                option_letter = match.group(1).upper()
                option_text = match.group(2).strip()
                options[option_letter] = option_text
            
            # If no multiline options found, try single-line format
            if len(options) == 0:
                # Pattern for single-line format: A) text B) text C) text D) text
                # More flexible pattern that handles various separators
                singleline_option_pattern = r'([A-D])[\)\.]\s*([^A]*?)(?=\s*[A-D][\)\.]|\s*Answer[\s\:]|\s*Explanation[\s\:]|\Z)'
                singleline_matches = re.finditer(singleline_option_pattern, question_text, re.IGNORECASE)
                
                for match in singleline_matches:
                    option_letter = match.group(1).upper()
                    option_text = match.group(2).strip()
                    # Clean up any trailing text that's not part of the option
                    option_text = re.split(r'\s*(?:Answer|Explanation)[\s\:]', option_text, flags=re.IGNORECASE)[0].strip()
                    # Remove any remaining "Answer" or "Explanation" fragments
                    option_text = re.sub(r'\s*(?:Answer|Explanation).*$', '', option_text, flags=re.IGNORECASE).strip()
                    options[option_letter] = option_text
            
            # Verify we have all 4 options
            if len(options) != 4 or set(options.keys()) != {'A', 'B', 'C', 'D'}:
                return None
            
            # Extract answer - more flexible pattern to handle "Answer: C" format
            answer_match = re.search(r'Answer[\s\:\.]*\s*([A-D])', question_text, re.IGNORECASE)
            if not answer_match:
                return None
            
            correct_answer = answer_match.group(1).upper()
            
            # Extract explanation (optional) - more flexible
            explanation_match = re.search(r'Explanation[\:\.]?\s*(.+?)(?=\n(?:Question|Q)\s*\d+|\n\d+\.|\Z)', 
                                         question_text, re.DOTALL | re.IGNORECASE)
            explanation = explanation_match.group(1).strip() if explanation_match else ""
            
            # Create question dictionary
            question_dict = {
                'question_number': question_number,
                'question_text': FileParserService.sanitize_question_text(question_text_content),
                'options': {
                    'A': FileParserService.sanitize_question_text(options['A']),
                    'B': FileParserService.sanitize_question_text(options['B']),
                    'C': FileParserService.sanitize_question_text(options['C']),
                    'D': FileParserService.sanitize_question_text(options['D'])
                },
                'correct_answer': correct_answer,
                'explanation': FileParserService.sanitize_question_text(explanation)
            }
            
            # Validate the parsed question
            if FileParserService.validate_question_format(question_dict):
                return question_dict
            
            return None
            
        except Exception as e:
            return None
    
    @staticmethod
    def parse_terms_conditions(file_path, file_type):
        """
        Parse Terms & Conditions from Word or PowerPoint file
        
        Args:
            file_path: Path to file
            file_type: 'word' or 'powerpoint'
            
        Returns:
            List of terms (max 10)
            
        Raises:
            ValidationError: If more than 10 bullets or invalid format
        """
        try:
            if file_type.lower() in ['word', 'docx']:
                terms = FileParserService._parse_terms_from_word(file_path)
            elif file_type.lower() in ['powerpoint', 'pptx']:
                terms = FileParserService._parse_terms_from_powerpoint(file_path)
            else:
                raise ValidationError("Invalid file type. Use 'word' or 'powerpoint'")
            
            # Validate maximum 10 bullets
            if len(terms) > 10:
                raise ValidationError(f"Too many terms. Maximum 10 allowed, found {len(terms)}")
            
            if not terms:
                raise ValidationError("No terms found in the document")
            
            return terms
            
        except Exception as e:
            if isinstance(e, ValidationError):
                raise
            raise ValidationError(f"Error parsing terms & conditions: {str(e)}")
    
    @staticmethod
    def _parse_terms_from_word(file_path):
        """Parse terms from Word document"""
        doc = Document(file_path)
        terms = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Remove numbering (1., 2., etc.) or bullet points
            cleaned_text = re.sub(r'^\d+[\.\)]\s*', '', text)
            cleaned_text = re.sub(r'^[\•\-\*]\s*', '', cleaned_text)
            cleaned_text = cleaned_text.strip()
            
            if cleaned_text:
                terms.append(cleaned_text)
        
        return terms[:10]  # Return only first 10
    
    @staticmethod
    def _parse_terms_from_powerpoint(file_path):
        """Parse terms from PowerPoint"""
        prs = Presentation(file_path)
        terms = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Split by newlines to get individual terms
                    lines = text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                        
                        # Remove numbering or bullet points
                        cleaned_line = re.sub(r'^\d+[\.\)]\s*', '', line)
                        cleaned_line = re.sub(r'^[\•\-\*]\s*', '', cleaned_line)
                        cleaned_line = cleaned_line.strip()
                        
                        if cleaned_line:
                            terms.append(cleaned_line)
        
        return terms[:10]  # Return only first 10
    
    @staticmethod
    def validate_question_format(question_dict):
        """
        Validate that question dictionary has all required fields
        
        Args:
            question_dict: Question dictionary to validate
            
        Returns:
            bool: True if valid, False otherwise
        """
        if not isinstance(question_dict, dict):
            return False
        
        # Check required fields
        required_fields = ['question_number', 'question_text', 'options', 'correct_answer']
        for field in required_fields:
            if field not in question_dict or not question_dict[field]:
                return False
        
        # Check question text is not empty
        if not question_dict['question_text'].strip():
            return False
        
        # Check options
        options = question_dict.get('options', {})
        if not isinstance(options, dict):
            return False
        
        # Check all options exist and are not empty
        for letter in ['A', 'B', 'C', 'D']:
            if letter not in options or not options[letter].strip():
                return False
        
        # Check correct answer is valid
        correct_answer = question_dict.get('correct_answer', '').upper()
        if correct_answer not in ['A', 'B', 'C', 'D']:
            return False
        
        return True
    
    @staticmethod
    def sanitize_question_text(text):
        """
        Clean and sanitize question text
        
        Args:
            text: Text to sanitize
            
        Returns:
            str: Sanitized text
        """
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove special characters that may break encryption (keep basic punctuation)
        # Allow: letters, numbers, spaces, and common punctuation
        text = re.sub(r'[^\w\s\.\,\?\!\-\:\;\(\)\'\"]', '', text)
        
        # Trim whitespace
        text = text.strip()
        
        return text
